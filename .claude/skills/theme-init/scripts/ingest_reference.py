#!/usr/bin/env python3
"""Extract layout-device signals from a reference deck (SVG or PPTX).

Mirrors slide-html's ingest_reference.py device set + JSON schema +
recommend-don't-apply contract, but parses slide-svg's input formats (SVG,
PPTX) instead of HTML. It is a RECOMMENDER: it never edits theme-active.json,
never re-renders, and never replaces the human-authored DESIGN.md §5/§6 — it
only adds a clearly-labeled, marker-fenced hint block. See
references/reference-ingestion.md for the output JSON schema.

Usage:
    python3 ingest_reference.py --reference <deck-dir|.svg|.pptx> \\
        [--out devices.json] [--design-md <DESIGN.md>] [--no-seed] [--quiet]
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

SCRIPTS = Path(__file__).resolve().parent
# .../.claude/skills/theme-init/scripts -> parents[1] == .../.claude/skills
THEME_ACTIVE = SCRIPTS.parents[1] / "slide" / "references" / "theme-active.json"

SCHEMA_VERSION = "1.0"

_RECT_RE = re.compile(r"<rect\b([^>]*?)/?>", re.I)
_LINE_RE = re.compile(r"<line\b([^>]*?)/?>", re.I)
_TEXT_RE = re.compile(r"<text\b([^>]*)>(.*?)</text>", re.I | re.S)
_ATTR_RE = re.compile(r"""([\w:-]+)\s*=\s*(["'])(.*?)\2""")
_TAG_RE = re.compile(r"<[^>]+>")


def _attrs(s: str) -> dict[str, str]:
    return {m.group(1).lower(): m.group(3) for m in _ATTR_RE.finditer(s)}


_NUM_RE = re.compile(r"\s*(-?\d*\.?\d+)")


def _num(v: Any, default: float = 0.0) -> float:
    if v is None:
        return default
    m = _NUM_RE.match(str(v))
    if not m:
        return default
    try:
        return float(m.group(1))
    except ValueError:  # pathological capture (e.g. lone '.') — degrade, never crash
        return default


def _has(attr: dict, key: str) -> bool:
    v = attr.get(key, "")
    return bool(v) and v.strip().lower() != "none"


def _is_kicker(attr: dict, text: str) -> bool:
    t = _TAG_RE.sub("", text).strip()
    if not t:
        return False
    fs = _num(attr.get("font-size"), 99)
    ls = _num(attr.get("letter-spacing"), 0)
    alpha = [c for c in t if c.isalpha()]
    upper = bool(alpha) and sum(c.isupper() for c in alpha) >= 0.8 * len(alpha)
    return fs <= 14 and (upper or ls > 0)


def _svg_canvas(svg: str) -> tuple[float, float]:
    """Canvas (width, height) from the <svg> viewBox, else its width/height attrs,
    else the 1280x720 default. Sizes the full-bleed-background skip so reference
    decks authored at other sizes (e.g. 1920x1080) are read correctly."""
    m = re.search(r'viewBox\s*=\s*["\']\s*[-\d.]+\s+[-\d.]+\s+([\d.]+)\s+([\d.]+)', svg, re.I)
    if m:
        return _num(m.group(1), 1280), _num(m.group(2), 720)
    mw = re.search(r'<svg\b[^>]*?\bwidth\s*=\s*["\']([\d.]+)', svg, re.I)
    mh = re.search(r'<svg\b[^>]*?\bheight\s*=\s*["\']([\d.]+)', svg, re.I)
    return (_num(mw.group(1), 1280) if mw else 1280.0,
            _num(mh.group(1), 720) if mh else 720.0)


def _analyze_svg(stem: str, svg: str) -> dict[str, Any]:
    filled = hairline = cta = 0
    fills: set[str] = set()
    rules = len(_LINE_RE.findall(svg))
    cw, ch = _svg_canvas(svg)
    for m in _RECT_RE.finditer(svg):
        a = _attrs(m.group(1))
        w, h = _num(a.get("width")), _num(a.get("height"))
        if w >= 0.93 * cw and h >= 0.93 * ch:   # full-canvas background
            continue
        if (0 < h <= 3 and w > h) or (0 < w <= 3 and h > w):  # thin rule
            rules += 1
            continue
        rounded = "rx" in a or "ry" in a
        has_stroke = _has(a, "stroke")   # a stroke COLOR makes a card hairline;
        #          a leftover stroke-width on stroke="none" must NOT (Figma/PPT export)
        has_fill = _has(a, "fill")
        if rounded and (has_stroke or has_fill):
            if has_stroke:
                hairline += 1
            else:
                filled += 1
                fills.add(a.get("fill", "").lower())
            if 80 <= w <= 360 and 28 <= h <= 72:
                cta += 1
        elif has_fill:
            fills.add(a.get("fill", "").lower())
    kicker = sum(1 for m in _TEXT_RE.finditer(svg)
                 if _is_kicker(_attrs(m.group(1)), m.group(2)))
    return {"slide": stem, "cards": filled + hairline, "filled": filled,
            "hairline": hairline, "surface_alt": len(fills) >= 2,
            "rules": rules, "cta": cta, "kicker": kicker}


def _pptx_fill_color(shape) -> str | None:
    try:
        from pptx.enum.dml import MSO_FILL_TYPE
        f = shape.fill
        if f.type == MSO_FILL_TYPE.SOLID:
            return str(f.fore_color.rgb)
    except Exception:  # noqa: BLE001 — shapes without a fill raise; treat as none
        pass
    return None


def _pptx_has_line(shape) -> bool:
    try:
        w = shape.line.width
        return w is not None and int(w) > 0
    except Exception:  # noqa: BLE001
        return False


def _analyze_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
    except ImportError:  # older python-pptx
        MSO_SHAPE = None
    prs = Presentation(str(path))
    sw = prs.slide_width or 12192000
    sh = prs.slide_height or 6858000
    rows: list[dict] = []
    for i, slide in enumerate(prs.slides, 1):
        filled = hairline = rules = cta = kicker = 0
        fills: set[str] = set()
        for shape in slide.shapes:
            w = int(shape.width or 0)
            h = int(shape.height or 0)
            is_line = False
            try:
                is_line = shape.shape_type == MSO_SHAPE_TYPE.LINE
            except Exception:  # noqa: BLE001
                pass
            if is_line or (0 < h <= sh * 0.012) or (0 < w <= sw * 0.012):
                rules += 1
                continue
            if w >= 0.95 * sw and h >= 0.95 * sh:        # full-bleed background
                fc = _pptx_fill_color(shape)
                if fc:
                    fills.add(fc)
                continue
            is_rect = False
            try:
                if MSO_SHAPE is not None:
                    is_rect = shape.auto_shape_type in (MSO_SHAPE.ROUNDED_RECTANGLE,
                                                        MSO_SHAPE.RECTANGLE)
            except Exception:  # noqa: BLE001 — non-autoshapes raise
                pass
            if not is_rect:
                continue
            if _pptx_has_line(shape):
                hairline += 1
            else:
                fc = _pptx_fill_color(shape)
                if fc:
                    filled += 1
                    fills.add(fc)
            try:
                txt = shape.text_frame.text.strip() if shape.has_text_frame else ""
            except Exception:  # noqa: BLE001
                txt = ""
            if txt and 0 < w <= sw * 0.3 and 0 < h <= sh * 0.12:
                cta += 1
            if txt and txt.isupper() and len(txt) <= 40:
                kicker += 1
        rows.append({"slide": f"slide{i:02d}", "cards": filled + hairline,
                     "filled": filled, "hairline": hairline,
                     "surface_alt": len(fills) >= 2, "rules": rules,
                     "cta": cta, "kicker": kicker})
    return rows


def _decide_card_style(filled: int, hairline: int, cards: int,
                       rules: int) -> tuple[str | None, str, str]:
    if cards == 0:
        if rules > 0:
            return ("borderless", "medium",
                    f"no card shapes; {rules} hairline rules (grouped by whitespace)")
        return None, "low", "no card-like shapes and no rules detected"
    if hairline >= filled:
        conf = "high" if hairline >= 2 * max(filled, 1) else "medium"
        return "hairline", conf, f"{hairline} bordered vs {filled} filled card shapes"
    conf = "high" if filled >= 2 * max(hairline, 1) else "medium"
    return "filled", conf, f"{filled} filled vs {hairline} bordered card shapes"


def _empty_devices() -> dict[str, Any]:
    return {
        "card_style": {"value": None, "confidence": "low", "evidence": "not analyzed"},
        "surface_alternation": {"present": False, "slides_using": 0, "ratio": 0.0},
        "hairline_dividers": {"present": False, "count": 0},
        "cta": {"present": False, "count": 0},
        "kicker": {"present": False, "count": 0},
    }


def _aggregate(ref: Path, fmt: str, per_slide: list[dict]) -> dict[str, Any]:
    n = len(per_slide)
    filled = sum(r["filled"] for r in per_slide)
    hairline = sum(r["hairline"] for r in per_slide)
    cards = sum(r["cards"] for r in per_slide)
    rules = sum(r["rules"] for r in per_slide)
    cta = sum(r["cta"] for r in per_slide)
    kicker = sum(r["kicker"] for r in per_slide)
    surf = sum(1 for r in per_slide if r["surface_alt"])
    cs_val, cs_conf, cs_ev = _decide_card_style(filled, hairline, cards, rules)
    public = [{"slide": r["slide"], "cards": r["cards"], "surface_alt": r["surface_alt"],
               "rules": r["rules"], "cta": r["cta"], "kicker": r["kicker"]}
              for r in per_slide]
    return {
        "version": SCHEMA_VERSION, "source": str(ref), "source_format": fmt,
        "slides_analyzed": n,
        "devices": {
            "card_style": {"value": cs_val, "confidence": cs_conf, "evidence": cs_ev},
            "surface_alternation": {"present": surf > 0, "slides_using": surf,
                                    "ratio": round(surf / n, 2) if n else 0.0},
            "hairline_dividers": {"present": rules > 0, "count": rules},
            "cta": {"present": cta > 0, "count": cta},
            "kicker": {"present": kicker > 0, "count": kicker},
        },
        "per_slide": public,
    }


_REF_START = "<!-- REF-INGEST:start -->"
_REF_END = "<!-- REF-INGEST:end -->"


def _seed_design(design_path: Path, result: dict) -> bool:
    """Insert/replace an additive, marker-fenced reference-hint block in DESIGN.md
    §5. Never touches the section's AGENT-FILL placeholders. Idempotent."""
    if not design_path.exists():
        return False
    d = result["devices"]
    cs = d["card_style"]
    block = (
        f"{_REF_START}\n"
        f"> **참조 자동추출 초안** — `{result['source']}` "
        f"({result['slides_analyzed']} slides, {result['source_format']}) 에서 측정한 레이아웃 신호. "
        f"힌트일 뿐이며 정식 §5/§6 항목은 직접 검토·작성해야 합니다.\n"
        f"> - card_style: **{cs['value']}** [{cs['confidence']}]\n"
        f"> - surface 교차: {d['surface_alternation']['slides_using']}/{result['slides_analyzed']} slides\n"
        f"> - hairline divider: {d['hairline_dividers']['count']} · "
        f"CTA: {d['cta']['count']} · kicker: {d['kicker']['count']}\n"
        f"{_REF_END}"
    )
    text = design_path.read_text(encoding="utf-8")
    if _REF_START in text and _REF_END in text:
        # lambda replacement: `block` holds a Windows path (C:\Users\...), so it
        # must NOT go through re.sub's backslash-escape processing.
        text = re.sub(re.escape(_REF_START) + r".*?" + re.escape(_REF_END),
                      lambda _m: block, text, flags=re.DOTALL)
    else:
        m5 = re.search(r"^## 5\..*$", text, re.MULTILINE)
        if m5:
            text = text[:m5.end()] + "\n\n" + block + text[m5.end():]
        else:
            text = text.rstrip() + "\n\n" + block + "\n"
    design_path.write_text(text, encoding="utf-8")
    return True


def summarize(result: dict) -> str:
    d = result["devices"]
    cs = d["card_style"]
    return "\n".join([
        f"reference: {result['source']}  ({result['slides_analyzed']} slides, {result['source_format']})",
        f"  card_style    : {cs['value']}  ({cs['confidence']} - {cs['evidence']})",
        f"  surface alt   : {d['surface_alternation']['slides_using']} slides "
        f"({int(d['surface_alternation']['ratio'] * 100)}%)",
        f"  hairline rules: {d['hairline_dividers']['count']}",
        f"  CTA           : {d['cta']['count']}  (low confidence for vector formats)",
        f"  kicker/eyebrow: {d['kicker']['count']}  (low confidence for vector formats)",
    ])


def _find_slides(ref: Path) -> tuple[list[Path], str | None]:
    """(slide files, source_format). SVG preferred over PPTX. None == unsupported."""
    if ref.is_file():
        suf = ref.suffix.lower()
        if suf == ".svg":
            return [ref], "svg"
        if suf == ".pptx":
            return [ref], "pptx"
        return [], None
    for sub in ("svg_final", "svg_output"):
        d = ref / sub
        if d.is_dir():
            svgs = sorted(p for p in d.glob("*.svg"))
            if svgs:
                return svgs, "svg"
    svgs = sorted(ref.glob("*.svg"))
    if svgs:
        return svgs, "svg"
    pptx = sorted(ref.glob("*.pptx"))
    if pptx:
        return [pptx[-1]], "pptx"
    return [], None


def extract(ref: Path) -> dict[str, Any]:
    slide_files, fmt = _find_slides(ref)
    if fmt is None:
        return {"version": SCHEMA_VERSION, "source": str(ref), "source_format": None,
                "slides_analyzed": 0, "devices": _empty_devices(), "per_slide": []}
    if fmt == "svg":
        per_slide = [_analyze_svg(p.stem, p.read_text(encoding="utf-8", errors="replace"))
                     for p in slide_files]
    else:
        try:
            per_slide = _analyze_pptx(slide_files[0])
        except Exception:  # noqa: BLE001 — corrupt .pptx / missing python-pptx →
            per_slide = []  # degrade to 0 slides so main() exits 1 with guidance
    return _aggregate(ref, fmt, per_slide)


def _current_card_style() -> str:
    try:
        j = json.loads(THEME_ACTIVE.read_text(encoding="utf-8"))
        return (j.get("surface") or {}).get("card_style") or "hairline"
    except (OSError, json.JSONDecodeError):
        return "hairline"


def main() -> int:
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8")  # cp949-console safety (Windows)
        except (AttributeError, ValueError):
            pass
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--reference", type=Path, required=True,
                    help="Reference deck dir, or a single .svg / .pptx file")
    ap.add_argument("--out", type=Path, default=None, help="Write JSON here (default: stdout)")
    ap.add_argument("--design-md", type=Path, default=None,
                    help="Rendered per-theme DESIGN.md to seed the §5 hint block into")
    ap.add_argument("--no-seed", action="store_true", help="Do not touch DESIGN.md")
    ap.add_argument("--quiet", action="store_true", help="JSON only, no stderr summary")
    args = ap.parse_args()

    if not args.reference.exists():
        print(f"reference not found: {args.reference}", file=sys.stderr)
        return 2
    result = extract(args.reference)
    if result["source_format"] is None:
        print(f"unsupported reference format: {args.reference}\n"
              f"  export the deck to PPTX or SVG, or fill DESIGN.md §5/§6 by hand.",
              file=sys.stderr)
        return 1
    if result["slides_analyzed"] == 0:
        print(f"no SVG/PPTX slides found under {args.reference}", file=sys.stderr)
        return 1

    payload = json.dumps(result, indent=2, ensure_ascii=False)
    if args.out:
        args.out.write_text(payload + "\n", encoding="utf-8")
        print(f"wrote {args.out}", file=sys.stderr)
    else:
        print(payload)

    rec = result["devices"]["card_style"]["value"]
    cur = _current_card_style()
    if rec and rec != cur:
        print(f"\n[!] reference suggests surface.card_style='{rec}' (active is '{cur}').\n"
              f"    to apply: set surface.card_style='{rec}' in the Step 2 theme JSON, "
              f"then re-render (init_theme.py). This script never edits the theme.",
              file=sys.stderr)
    elif rec:
        print(f"\n[ok] reference card_style '{rec}' matches the active theme.", file=sys.stderr)

    if not args.no_seed and args.design_md:
        if _seed_design(args.design_md, result):
            print("  - seeded DESIGN.md §5 reference-hint block (additive; placeholders intact)",
                  file=sys.stderr)

    if not args.quiet:
        print("\n" + summarize(result), file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
